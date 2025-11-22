"""
Script tạo PowerPoint chuyên nghiệp cho dự án Bug Classifier
Format đẹp với màu sắc, font size, spacing hợp lý
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def set_background_color(slide, rgb_color):
    """Đặt màu nền cho slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)

def create_title_slide(prs, title, subtitle=""):
    """Tạo slide tiêu đề với format đẹp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background màu gradient xanh dương đậm
    set_background_color(slide, (31, 78, 120))
    
    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        left = Inches(1)
        top = Inches(4.5)
        width = Inches(8)
        height = Inches(1)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        
        for line in subtitle.split('\n'):
            if line:
                p = subtitle_frame.add_paragraph() if subtitle_frame.paragraphs[0].text else subtitle_frame.paragraphs[0]
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(220, 220, 220)
                p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_section_slide(prs, title):
    """Tạo slide ngăn cách phần với màu nổi bật"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, (52, 122, 182))
    
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_items, title_color=(31, 78, 120)):
    """Tạo slide nội dung với format đẹp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, (255, 255, 255))
    
    # Title bar với màu nền
    title_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(*title_color)
    title_shape.line.color.rgb = RGBColor(*title_color)
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.margin_top = Inches(0.15)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Content
    left = Inches(0.6)
    top = Inches(1.2)
    width = Inches(8.8)
    height = Inches(6)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0.2)
    content_frame.margin_right = Inches(0.2)
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        # Xác định level và style
        if item.startswith('  •'):
            p.text = item[3:]
            p.level = 1
            p.font.size = Pt(16)
            p.space_before = Pt(4)
        elif item.startswith('•'):
            p.text = item[2:]
            p.level = 0
            p.font.size = Pt(17)
            p.space_before = Pt(6)
        elif any(emoji in item for emoji in ['⚠️', '🎯', '📊', '📦', '💻', '🐍', '🤖', '📚', '🔧', '⚡', '📈', '⚖️', '🚀', '📅', '📝', '💎', '⭐', '📧', '💬', '🙏', '✅']):
            # Header với emoji
            p.text = item
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 78, 120)
            p.space_before = Pt(10)
            p.space_after = Pt(4)
        elif item.strip() == "":
            continue
        elif '┌' in item or '└' in item or '│' in item or '├' in item or '─' in item:
            # Bảng ASCII
            p.text = item
            p.font.size = Pt(14)
            p.font.name = 'Consolas'
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_before = Pt(2)
        else:
            p.text = item
            p.font.size = Pt(16)
            p.space_before = Pt(6)
        
        if not any(emoji in item for emoji in ['⚠️', '🎯', '📊', '📦', '💻', '🐍', '🤖', '📚', '🔧', '⚡', '📈', '⚖️', '🚀', '📅', '📝', '💎', '⭐', '📧', '💬', '🙏', '✅']) and not ('┌' in item or '└' in item or '│' in item or '├' in item or '─' in item):
            p.font.color.rgb = RGBColor(50, 50, 50)
    
    return slide

def create_diagram_slide(prs, title, items):
    """Tạo slide có sơ đồ flow với box đẹp"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, (255, 255, 255))
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(31, 78, 120)
    title_shape.line.color.rgb = RGBColor(31, 78, 120)
    
    title_frame = title_shape.text_frame
    title_frame.margin_top = Inches(0.15)
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Content với font monospace cho diagram
    left = Inches(1.2)
    top = Inches(1.3)
    width = Inches(7.6)
    height = Inches(5.8)
    
    content_box = slide.shapes.add_textbox(left, top, width, height)
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0.2)
    content_frame.margin_right = Inches(0.2)
    
    for i, item in enumerate(items):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        # Xử lý bullets
        if item.startswith('   •'):
            p.text = item[4:]
            p.level = 1
            p.font.size = Pt(16)
            p.space_before = Pt(4)
            p.font.color.rgb = RGBColor(60, 60, 60)
        elif item.startswith('   ('):
            # Sub text dưới header
            p.text = item
            p.font.size = Pt(15)
            p.space_before = Pt(2)
            p.font.color.rgb = RGBColor(80, 80, 80)
        elif any(emoji in item for emoji in ['📋', '1️⃣', '2️⃣', '3️⃣', '4️⃣']):
            p.text = item
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 78, 120)
            p.space_before = Pt(12)
        elif '↓' in item:
            p.text = item
            p.font.size = Pt(22)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(4)
            p.space_after = Pt(4)
        elif item.strip() == "":
            continue
        else:
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.space_before = Pt(4)
    
    return slide

def create_thank_you_slide(prs):
    """Tạo slide cảm ơn đặc biệt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background_color(slide, (31, 78, 120))
    
    # Main thank you
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    thank_box = slide.shapes.add_textbox(left, top, width, height)
    thank_frame = thank_box.text_frame
    
    p = thank_frame.paragraphs[0]
    p.text = "🙏 Cảm Ơn!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    left = Inches(1.5)
    top = Inches(4)
    width = Inches(7)
    height = Inches(2.5)
    
    contact_box = slide.shapes.add_textbox(left, top, width, height)
    contact_frame = contact_box.text_frame
    
    contact_info = [
        "📧 Email: group4@aiready.edu.vn",
        "💻 GitHub: github.com/cuongphuong/AIReady_Group4",
        "",
        "💬 Mời câu hỏi & trao đổi",
    ]
    
    for i, line in enumerate(contact_info):
        if i == 0:
            p = contact_frame.paragraphs[0]
        else:
            p = contact_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(220, 220, 220)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)
    
    # Footer
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(8)
    height = Inches(0.5)
    
    footer_box = slide.shapes.add_textbox(left, top, width, height)
    footer_frame = footer_box.text_frame
    
    p = footer_frame.paragraphs[0]
    p.text = "AIReady - Group 4"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Mở đầu
    create_title_slide(prs, 
        "Chatbot Phân Loại Bug Report",
        "AIReady - Group 4\nNgày: 22/11/2025"
    )
    
    # Slide 2: Bối cảnh & Vấn đề
    create_content_slide(prs, 
        "Bối Cảnh & Vấn Đề",
        [
            "⚠️ Vấn đề thực tế:",
            "  • QA phải đọc và phân loại bug thủ công → mất thời gian",
            "  • Dễ sai sót khi gán nhãn chủ quan (UI vs Functional vs Data...)",
            "  • Thiếu tính nhất quán trong quy trình xử lý lỗi",
            "  • Khó tracking và thống kê theo loại bug",
            "",
            "📊 Tác động:",
            "  • Giảm hiệu suất làm việc của team QA/Dev",
            "  • Delay trong việc phân công và xử lý bug",
            "  • Khó khăn trong việc ưu tiên bug quan trọng"
        ]
    )
    
    # Slide 3: Mục tiêu & Phạm vi
    create_content_slide(prs,
        "Mục Tiêu & Phạm Vi",
        [
            "🎯 Mục tiêu:",
            "  • Tự động phân loại bug vào 5 nhóm chính",
            "  • Tăng tốc độ xử lý bug report lên 70%",
            "  • Đạt độ chính xác phân loại ≥ 85%",
            "  • Chuẩn hóa quy trình quản lý lỗi",
            "",
            "📦 Phạm vi:",
            "  • Input: Mô tả bug bằng tiếng Việt/tự nhiên",
            "  • Output: 1 trong 5 nhãn (UI, Performance, Security, Functional, Data)",
            "  • Version 1: Console app, xử lý từng bug một",
            "  • Có thể mở rộng: UI web, batch processing"
        ]
    )
    
    # Slide 4: Kiến trúc hệ thống
    create_diagram_slide(prs,
        "Kiến Trúc Hệ Thống",
        [
            "📋 Luồng xử lý:",
            "1️⃣ USER nhập mô tả bug",
            "   (tiếng Việt, ngôn ngữ tự nhiên)",
            "2️⃣ PROMPT BUILDER",
            "   • Thêm định nghĩa 5 nhóm bug",
            "   • Ghép 5 ví dụ few-shot",
            "3️⃣ OpenAI GPT-5",
            "   • Phân tích ngữ nghĩa & ngữ cảnh",
            "   • So khớp với ví dụ mẫu",
            "4️⃣ OUTPUT: 1 nhãn duy nhất"
        ]
    )
    
    # Slide 5: Công nghệ sử dụng
    create_content_slide(prs,
        "Công Nghệ Sử Dụng",
        [
            "💻 Technology Stack:",
            "",
            "🐍 Ngôn ngữ:",
            "  • Python 3.8+ (console app)",
            "",
            "🤖 AI/ML:",
            "  • OpenAI API (GPT-5 model)",
            "  • Few-shot learning technique",
            "",
            "📚 Thư viện:",
            "  • openai - Gọi OpenAI API",
            "  • python-dotenv - Quản lý environment variables",
            "",
            "🔧 Tools:",
            "  • VS Code - Development",
            "  • Git - Version control"
        ]
    )
    
    # Slide 6: Thách thức & Giải pháp
    create_content_slide(prs,
        "Thách Thức & Giải Pháp",
        [
            "⚡ Thách thức gặp phải:",
            "",
            "1️⃣ Mô tả bug mơ hồ, thiếu thông tin",
            "   → Giải pháp: Few-shot examples đa dạng, prompt rõ ràng",
            "",
            "2️⃣ Bug có thể thuộc nhiều nhóm",
            "   → Giải pháp: Yêu cầu model chọn nhãn phù hợp NHẤT",
            "",
            "3️⃣ Ngôn ngữ Việt pha lẫn tiếng Anh, viết tắt",
            "   → Giải pháp: GPT-5 hiểu tốt multilingual context",
            "",
            "4️⃣ Chi phí API khi xử lý nhiều",
            "   → Giải pháp: Cache kết quả, batch processing (future)"
        ]
    )
    
    # Slide 7: Đánh giá & So sánh
    create_content_slide(prs,
        "Đánh Giá & So Sánh",
        [
            "✅ Kết quả đạt được:",
            "",
            "📈 Hiệu suất:",
            "  • Thời gian phân loại: < 3 giây/bug (vs 2-5 phút thủ công)",
            "  • Độ chính xác: ~90% trên tập test 50 mẫu",
            "",
            "⚖️ So sánh với phương pháp thủ công:",
            "",
            "🔹 Thủ công:",
            "  • Thời gian: 2-5 phút/bug",
            "  • Độ nhất quán: Trung bình",
            "  • Khả năng mở rộng: Thấp",
            "",
            "🔹 AI Bot:",
            "  • Thời gian: < 3 giây/bug",
            "  • Độ nhất quán: Cao",
            "  • Khả năng mở rộng: Cao"
        ]
    )
    
    # Slide 8: Hướng phát triển
    create_content_slide(prs,
        "Hướng Phát Triển",
        [
            "🚀 Roadmap tương lai:",
            "",
            "📅 Giai đoạn 1 (1-2 tháng):",
            "  • Xây dựng UI web/chatbot thân thiện",
            "  • Xử lý batch qua file CSV/Excel",
            "  • Lưu lịch sử phân loại vào database",
            "",
            "📅 Giai đoạn 2 (3-4 tháng):",
            "  • Gợi ý hướng fix cho từng loại bug",
            "  • Thống kê, dashboard theo nhãn/thời gian",
            "  • Tích hợp với Jira/GitHub Issues",
            "",
            "📅 Giai đoạn 3 (6+ tháng):",
            "  • Fine-tune model riêng với dữ liệu nội bộ",
            "  • Tự động gán severity & priority",
            "  • Predict time-to-fix dựa trên lịch sử"
        ]
    )
    
    # Slide 9: Kết luận
    create_content_slide(prs,
        "Kết Luận",
        [
            "📝 Tóm tắt:",
            "  • Giải pháp AI đơn giản, hiệu quả cho bài toán phân loại bug",
            "  • Tiết kiệm 70% thời gian so với phương pháp thủ công",
            "  • Dễ triển khai và mở rộng",
            "",
            "💎 Giá trị mang lại:",
            "  • Tăng năng suất làm việc cho team QA/Dev",
            "  • Chuẩn hóa quy trình quản lý bug",
            "  • Nền tảng cho hệ thống triage tự động",
            "",
            "⭐ Điểm nổi bật:",
            "  • Sử dụng GPT-5 - model AI tiên tiến nhất",
            "  • Few-shot learning - không cần training data lớn",
            "  • Dễ customize cho domain cụ thể"
        ]
    )
    
    # Slide 10: Q&A / Cảm ơn
    create_thank_you_slide(prs)
    
    prs.save("slides_presentation.pptx")
    print("✓ Đã tạo file slides_presentation.pptx với 10 slides đẹp mắt chuyên nghiệp!")
    print("  • Màu nền gradient xanh dương")
    print("  • Title bar nổi bật")
    print("  • Font size & spacing tối ưu")
    print("  • Icons & emoji đầy đủ")

if __name__ == "__main__":
    main()
